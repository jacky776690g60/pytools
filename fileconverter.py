''' =================================================================
| fileconverter.py -- tools/fileconverter.py
|
| Created by Jack on 07/11/2023
| Copyright © 2023 jacktogon. All rights reserved.
================================================================= '''

import os, sys, argparse, time
from pathlib import Path
from io import BytesIO

PPT_TO_PDF_USABLE = True
try:
    from pdf2image import convert_from_path
    from pptx import Presentation
    from pptx.util import Inches
except ImportError:
    PPT_TO_PDF_USABLE = False



__all__ = ('FileConverter', )



class FileConverter():

    @staticmethod
    def pdf2ppt(
        pdf_path:       str, 
        high_res:       bool,
        img_extension:  str     = "png"
    ) -> None:
        '''
        Convert a single PDF file to a PowerPoint presentation.
        
        Requires Poppler to be installed first.

        Params:
        -------
        - `pdf_path`: Path to the PDF file.
        - `high_res`: If True, converts with higher DPI.
        - `img_extension`: img extension for the output. E.g., jpeg, png, tiff...
        '''
        if not PPT_TO_PDF_USABLE: raise ImportError(
                "The necessary packages are not installed to use this function.\n Use `pip install pdf2image python-pptx`"
            )
        
        ppt = Presentation() # Initialize a PowerPoint presentation
        BLANK_SLIDE_LAYOUT = ppt.slide_layouts[6]
        
        file_path = Path(pdf_path)
        base_name = file_path.stem
        
        # Convert PDF to images
        dpi = 150 if high_res else 75
        pdf_imgs = convert_from_path(str(file_path), dpi=dpi, fmt='ppm', thread_count=2)
        
        SCALER = 9525
        '''
        PowerPoint uses a different unit of measurement called "EMU" 
        (English Metric Units) for slide dimensions, whereas the dimensions 
        of the images obtained from the PDF are in pixels.
        
        1 EMU is defined as 1/914400 of an inch. To convert pixels to EMUs, 
        you typically need to know the DPI (dots per inch) of the image, and 
        then you can use the formula
        '''
        
        for img in pdf_imgs:
            img_buffer = BytesIO()          # in-memory byte buffer
            img.save(img_buffer, format = img_extension)
            img_buffer.seek(0)              # move pointer to start of buffer
            w, h = img.size
            
            ppt.slide_width  = w * SCALER
            ppt.slide_height = h * SCALER 
            # ◼︎ add to pptx
            slide = ppt.slides.add_slide(BLANK_SLIDE_LAYOUT)
            slide.shapes.add_picture(img_buffer, 0, 0, width=w * SCALER, height=h * SCALER)
        ppt.save(str(file_path.parent / base_name) + '.pptx')



if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="PDF to PPT converter tool")
    parser.add_argument("-p", "--pdf_path", type=str, required=True, metavar="",
                        help="path to your pdf file")
    parser.add_argument("-hs", "--high_res", type=bool, default=False, metavar="", action=argparse.BooleanOptionalAction,
                        help="Convert with higher dpi")
    args = parser.parse_args()

    FileConverter.pdf2ppt(args.pdf_path, args.high_res)

